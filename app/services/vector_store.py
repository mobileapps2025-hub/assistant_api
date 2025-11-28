import os
import json
import hashlib
import sys
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
import numpy as np

# Document processing imports
import PyPDF2
try:
    import fitz  # PyMuPDF
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# Semantic search imports
try:
    import faiss
    from sentence_transformers import SentenceTransformer
    HAS_VECTOR_DEPS = True
except ImportError:
    HAS_VECTOR_DEPS = False

from app.core.logging import get_logger

logger = get_logger(__name__)

class VectorStoreService:
    """
    Service for managing the persistent FAISS vector store.
    """
    
    INDEX_FILENAME = "faiss_index.bin"
    METADATA_FILENAME = "chunk_metadata.json"
    EMBEDDING_MODEL_NAME = 'all-MiniLM-L6-v2'

    def __init__(self, storage_path: str):
        """
        Initialize the VectorStoreService.

        Args:
            storage_path: Directory path where index files will be stored.
        """
        self.storage_path = Path(storage_path)
        self.index_path = self.storage_path / self.INDEX_FILENAME
        self.metadata_path = self.storage_path / self.METADATA_FILENAME
        
        # Ensure storage directory exists
        if not self.storage_path.exists():
            os.makedirs(self.storage_path, exist_ok=True)
            
        self.embedding_model = None
        self.index = None
        self.chunks = []

    def index_exists(self) -> bool:
        """
        Check if the vector index and metadata files exist on disk.

        Returns:
            True if both files exist, False otherwise.
        """
        return self.index_path.exists() and self.metadata_path.exists()

    def _initialize_model(self):
        """Initialize the SentenceTransformer model if not already loaded."""
        if not HAS_VECTOR_DEPS:
            raise ImportError("faiss or sentence_transformers not installed.")
            
        if self.embedding_model is None:
            logger.info(f"Loading embedding model: {self.EMBEDDING_MODEL_NAME}")
            self.embedding_model = SentenceTransformer(self.EMBEDDING_MODEL_NAME)

    def build_index(self, documents_path: str) -> Dict[str, Any]:
        """
        Build the vector index from documents in the specified path.
        
        Args:
            documents_path: Path to the directory containing documents.
            
        Returns:
            Dictionary with build statistics.
        """
        if not HAS_VECTOR_DEPS:
            return {"success": False, "error": "Missing dependencies (faiss/sentence-transformers)"}

        self._initialize_model()
        
        doc_path = Path(documents_path)
        if not doc_path.exists():
            return {"success": False, "error": f"Documents path not found: {documents_path}"}

        logger.info(f"Starting index build from: {doc_path}")
        
        # 1. Collect and Process Files
        all_chunks = []
        chunk_id_counter = 0
        processed_files = 0
        
        # Supported extensions
        extensions = ['.pdf', '.docx', '.pptx', '.md']
        files = []
        for ext in extensions:
            files.extend(list(doc_path.glob(f"*{ext}")))
            
        # Filter for MCL documents (using logic adapted from original services.py)
        mcl_files = [f for f in files if self._is_mcl_document(f)]
        
        for file_path in mcl_files:
            try:
                text = self._extract_text(file_path)
                if not text:
                    continue
                    
                file_chunks = self._create_chunks(text)
                
                for i, chunk_text in enumerate(file_chunks):
                    chunk_metadata = {
                        "chunk_id": chunk_id_counter,
                        "document_name": file_path.name,
                        "document_type": file_path.suffix.upper().replace('.', ''),
                        "chunk_index": i,
                        "total_chunks": len(file_chunks),
                        "content": chunk_text,
                        "content_hash": hashlib.md5(chunk_text.encode()).hexdigest()[:8],
                        "file_path": str(file_path)
                    }
                    all_chunks.append(chunk_metadata)
                    chunk_id_counter += 1
                
                processed_files += 1
                logger.info(f"Processed {file_path.name}: {len(file_chunks)} chunks")
                
            except Exception as e:
                logger.error(f"Error processing {file_path.name}: {e}")

        if not all_chunks:
            return {"success": False, "error": "No chunks created from documents."}

        # 2. Create Embeddings
        logger.info(f"Creating embeddings for {len(all_chunks)} chunks...")
        texts = [c['content'] for c in all_chunks]
        embeddings = self.embedding_model.encode(texts, show_progress_bar=True)
        
        # 3. Build FAISS Index
        dimension = embeddings.shape[1]
        # Use IndexFlatIP for cosine similarity (normalized vectors)
        index = faiss.IndexFlatIP(dimension)
        faiss.normalize_L2(embeddings)
        index.add(embeddings)
        
        self.index = index
        self.chunks = all_chunks
        
        # 4. Save to Disk
        self._save_to_disk()
        
        return {
            "success": True,
            "processed_files": processed_files,
            "total_chunks": len(all_chunks),
            "index_path": str(self.index_path)
        }

    def _save_to_disk(self):
        """Save the index and metadata to disk."""
        if self.index is None:
            return

        # Save FAISS index
        faiss.write_index(self.index, str(self.index_path))
        
        # Save metadata
        with open(self.metadata_path, 'w', encoding='utf-8') as f:
            json.dump(self.chunks, f, ensure_ascii=False, indent=2)
            
        logger.info(f"Index saved to {self.index_path}")
        logger.info(f"Metadata saved to {self.metadata_path}")

    def load_index(self) -> bool:
        """
        Load the vector index and metadata from disk.
        
        Returns:
            True if successful, False otherwise.
        """
        if not self.index_exists():
            logger.warning("Index files not found.")
            return False
            
        if not HAS_VECTOR_DEPS:
            logger.error("Missing dependencies (faiss/sentence-transformers).")
            return False

        try:
            logger.info(f"Loading index from {self.index_path}...")
            self.index = faiss.read_index(str(self.index_path))
            
            logger.info(f"Loading metadata from {self.metadata_path}...")
            with open(self.metadata_path, 'r', encoding='utf-8') as f:
                self.chunks = json.load(f)
                
            # Initialize model for future queries
            self._initialize_model()
            
            logger.info(f"Index loaded successfully with {len(self.chunks)} chunks.")
            return True
            
        except Exception as e:
            logger.error(f"Error loading index: {e}")
            return False

    def _is_mcl_document(self, file_path: Path) -> bool:
        """Determine if a document is MCL-related."""
        name = file_path.name.lower()
        # Exclude Spotplan
        if any(x in name for x in ['spotplan', 'spot_plan', 'spot-plan']):
            return False
        # Include MCL indicators
        indicators = [
            'mcl', 'checklist', 'creating', 'questions', 'knowledge', 'base',
            'rollenprofil', 'dashboard', 'tablet', 'phone', 'how-to-use',
            'aufgabe', 'mistakes', 'release', 'notes', 'quiz', 'dropbox',
            'business', 'case', 'tech_updates', 'vorgehen', 'app', 'tests'
        ]
        return any(x in name for x in indicators)

    def _extract_text(self, file_path: Path) -> str:
        """Extract text based on file extension."""
        ext = file_path.suffix.lower()
        if ext == '.pdf':
            return self._extract_from_pdf(file_path)
        elif ext == '.docx':
            return self._extract_from_docx(file_path)
        elif ext == '.pptx':
            return self._extract_from_pptx(file_path)
        elif ext == '.md':
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            except Exception:
                return ""
        return ""

    def _extract_from_pdf(self, file_path: Path) -> str:
        text = ""
        if HAS_PYMUPDF:
            try:
                with fitz.open(file_path) as doc:
                    for page in doc:
                        text += page.get_text() + "\n\n"
                return text.strip()
            except Exception:
                pass
        
        # Fallback
        try:
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() + "\n\n"
            return text.strip()
        except Exception:
            return ""

    def _extract_from_docx(self, file_path: Path) -> str:
        if not HAS_DOCX: return ""
        try:
            doc = Document(file_path)
            text = []
            for p in doc.paragraphs:
                if p.text.strip(): text.append(p.text.strip())
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip(): text.append(cell.text.strip())
            return "\n\n".join(text)
        except Exception:
            return ""

    def _extract_from_pptx(self, file_path: Path) -> str:
        if not HAS_PPTX: return ""
        try:
            pres = Presentation(file_path)
            text = []
            for slide in pres.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text.strip())
            return "\n".join(text)
        except Exception:
            return ""

    def _create_chunks(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """Split text into overlapping chunks."""
        if len(text) <= chunk_size:
            return [text]
        
        chunks = []
        start = 0
        while start < len(text):
            end = start + chunk_size
            if end < len(text):
                # Try to find a natural break point
                break_points = [
                    text.rfind('.', start, end),
                    text.rfind('\n', start, end),
                    text.rfind(' ', start, end)
                ]
                best_break = max([bp for bp in break_points if bp > start + chunk_size // 2], default=end)
                end = best_break + 1 if best_break != end else end
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            start = end - overlap if end < len(text) else end
            
        return chunks

    def search(self, query: str, limit: int = 5) -> List[Dict[str, Any]]:
        """
        Search the vector index for relevant chunks.
        
        Args:
            query: The search query string.
            limit: Maximum number of results to return.
            
        Returns:
            List of chunk metadata dictionaries with similarity scores.
        """
        if not HAS_VECTOR_DEPS:
            logger.error("Missing dependencies (faiss/sentence-transformers).")
            return []
            
        if self.index is None or not self.chunks:
            logger.warning("Index not loaded. Call load_index() first.")
            return []
            
        try:
            # Initialize model if needed (should be done in load_index, but safe check)
            self._initialize_model()
            
            # Create query embedding
            query_embedding = self.embedding_model.encode([query])
            faiss.normalize_L2(query_embedding)
            
            # Search index
            # k must be <= total chunks
            k = min(limit, len(self.chunks))
            scores, indices = self.index.search(query_embedding, k)
            
            results = []
            for score, idx in zip(scores[0], indices[0]):
                if idx < len(self.chunks):
                    chunk = self.chunks[idx].copy()
                    chunk['similarity_score'] = float(score)
                    results.append(chunk)
                    
            return results
            
        except Exception as e:
            logger.error(f"Error during search: {e}")
            return []
