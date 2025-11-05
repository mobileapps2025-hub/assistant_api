# -*- coding: utf-8 -*-
import json
import requests
from io import BytesIO
import os
import sys

# Configure stdout encoding for emoji support
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')

import PyPDF2
try:
    import fitz  # PyMuPDF
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False
    print("PyMuPDF not available, will use PyPDF2 only for PDF processing")

# Document processing imports
try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    print("python-docx not available, will skip DOCX processing")

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("python-pptx not available, will skip PPTX processing")

# Semantic search imports
try:
    from sentence_transformers import SentenceTransformer
    import faiss
    import numpy as np
    HAS_SEMANTIC_SEARCH = True
    print("‚úÖ Semantic search capabilities available")
except ImportError:
    HAS_SEMANTIC_SEARCH = False
    print("‚ö†Ô∏è Semantic search libraries not available, falling back to keyword search")

from pathlib import Path
from app.config import client, AsyncSessionLocal
from app.models import EventsBetweenWeeksRequest
# from app.utils import FUNCTION_TOOLS  # Commented out - utils.py removed
# from app.clients.api_client import APIClient  # Commented out - clients/ removed
from pydantic import BaseModel
from datetime import datetime
from typing import List, Dict, Any, Tuple
import hashlib
import asyncio
import traceback

# Global API client instance - will be set by the middleware for Spotplan
# _api_client: APIClient = None  # Commented out - APIClient removed

# Global variables for MCL knowledge base - SEPARATED FROM SPOTPLAN
_mcl_vector_store_id: str = None
_mcl_document_chunks: List[Dict[str, Any]] = []
_mcl_embeddings: np.ndarray = None
_mcl_faiss_index = None
_mcl_embedding_model = None

# Global variables for Spotplan knowledge base - SEPARATED FROM MCL  
_spotplan_vector_store_id: str = None
_spotplan_document_chunks: List[Dict[str, Any]] = []

# --- SPOTPLAN CODE REMOVED (clients/ and database/ directories removed) ---
# All Spotplan-specific functions have been deleted.
# This assistant now only handles MCL (My Checklist) document processing.

# --- Enhanced Document Processing Functions (SHARED BY MCL) ---

def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX file."""
    if not HAS_DOCX:
        print(f"DOCX processing not available, skipping: {file_path}")
        return ""
    
    try:
        doc = Document(file_path)
        text_content = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text.strip())
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text_content.append(cell.text.strip())
        
        return "\n\n".join(text_content)
    except Exception as e:
        print(f"Error extracting text from DOCX {file_path}: {e}")
        return ""

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX file."""
    if not HAS_PPTX:
        print(f"PPTX processing not available, skipping: {file_path}")
        return ""
    
    try:
        pres = Presentation(file_path)
        text_content = []
        
        for slide_num, slide in enumerate(pres.slides, 1):
            slide_text = [f"=== Slide {slide_num} ==="]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if len(slide_text) > 1:  # More than just the slide header
                text_content.extend(slide_text)
                text_content.append("")  # Add spacing between slides
        
        return "\n".join(text_content)
    except Exception as e:
        print(f"Error extracting text from PPTX {file_path}: {e}")
        return ""

def extract_text_from_file(file_path: Path) -> str:
    """Extract text from any supported file type."""
    file_extension = file_path.suffix.lower()
    
    print(f"üìÑ Processing {file_path.name} ({file_extension})")
    
    if file_extension == '.pdf':
        return extract_text_from_pdf(str(file_path))
    elif file_extension == '.docx':
        return extract_text_from_docx(str(file_path))
    elif file_extension == '.pptx':
        return extract_text_from_pptx(str(file_path))
    elif file_extension == '.md':
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            print(f"Error reading markdown file {file_path}: {e}")
            return ""
    else:
        print(f"‚ö†Ô∏è Unsupported file type: {file_extension}")
        return ""

def initialize_semantic_search():
    """Initialize semantic search model if available."""
    global _mcl_embedding_model
    
    if not HAS_SEMANTIC_SEARCH:
        print("‚ö†Ô∏è Semantic search not available, using keyword search")
        return False
    
    try:
        print("üîÑ Loading semantic search model...")
        _mcl_embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
        print("‚úÖ Semantic search model loaded successfully")
        return True
    except Exception as e:
        print(f"‚ùå Error loading semantic search model: {e}")
        return False

def create_embeddings_for_chunks(chunks: List[Dict[str, Any]]) -> Tuple[np.ndarray, Any]:
    """Create embeddings for document chunks and build FAISS index."""
    if not HAS_SEMANTIC_SEARCH or not _mcl_embedding_model:
        return None, None
    
    try:
        print(f"üîÑ Creating embeddings for {len(chunks)} chunks...")
        
        # Extract text content from chunks
        texts = [chunk['content'] for chunk in chunks]
        
        # Create embeddings
        embeddings = _mcl_embedding_model.encode(texts, show_progress_bar=True)
        
        # Create FAISS index
        dimension = embeddings.shape[1]
        index = faiss.IndexFlatIP(dimension)  # Inner product for cosine similarity
        
        # Normalize embeddings for cosine similarity
        faiss.normalize_L2(embeddings)
        index.add(embeddings)
        
        print(f"‚úÖ Created embeddings and FAISS index with {len(chunks)} vectors")
        return embeddings, index
        
    except Exception as e:
        print(f"‚ùå Error creating embeddings: {e}")
        return None, None

def semantic_search_chunks(query: str, chunks: List[Dict[str, Any]], 
                          embeddings: np.ndarray, faiss_index, max_chunks: int = 5) -> List[Dict[str, Any]]:
    """Perform semantic search using FAISS."""
    if not HAS_SEMANTIC_SEARCH or not _mcl_embedding_model or faiss_index is None:
        return []
    
    try:
        print(f"üîç Performing semantic search for: '{query}'")
        
        # Create query embedding
        query_embedding = _mcl_embedding_model.encode([query])
        faiss.normalize_L2(query_embedding)
        
        # Search
        scores, indices = faiss_index.search(query_embedding, min(max_chunks, len(chunks)))
        
        # Get results
        results = []
        for i, (score, idx) in enumerate(zip(scores[0], indices[0])):
            if idx < len(chunks):  # Valid index
                chunk = chunks[idx].copy()
                chunk['similarity_score'] = float(score)
                results.append(chunk)
                print(f"üéØ Semantic match {i+1}: {chunk['document_name']} (Score: {score:.3f})")
        
        return results
        
    except Exception as e:
        print(f"‚ùå Error in semantic search: {e}")
        return []

def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file using PyMuPDF for better text extraction."""
    text_content = ""
    
    # Try PyMuPDF first if available (better text extraction)
    if HAS_PYMUPDF:
        try:
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                page_text = page.get_text()
                if page_text:
                    text_content += page_text + "\n\n"
            doc.close()
            return text_content.strip()
        except Exception as e:
            print(f"Error extracting text from {file_path} with PyMuPDF: {e}")
    
    # Fallback to PyPDF2 (or use directly if PyMuPDF not available)
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text_content += page.extract_text() + "\n\n"
        return text_content.strip()
    except Exception as e2:
        print(f"Error extracting text from {file_path} with PyPDF2: {e2}")
        return ""

def create_text_chunks(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
    """Split text into overlapping chunks for better retrieval."""
    if len(text) <= chunk_size:
        return [text]
    
    chunks = []
    start = 0
    
    while start < len(text):
        end = start + chunk_size
        
        # Try to break at a sentence or paragraph boundary
        if end < len(text):
            # Look for the last period, newline, or space within the chunk
            break_points = [text.rfind('.', start, end), text.rfind('\n', start, end), text.rfind(' ', start, end)]
            best_break = max([bp for bp in break_points if bp > start + chunk_size // 2], default=end)
            end = best_break + 1 if best_break != end else end
        
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        
        start = end - overlap if end < len(text) else end
    
    return chunks

# --- MCL Knowledge Base Functions (SEPARATED FROM SPOTPLAN) ---

def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file using PyMuPDF for better text extraction."""
    text_content = ""
    
    # Try PyMuPDF first if available (better text extraction)
    if HAS_PYMUPDF:
        try:
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                page_text = page.get_text()
                if page_text:
                    text_content += page_text + "\n\n"
            doc.close()
            return text_content.strip()
        except Exception as e:
            print(f"Error extracting text from {file_path} with PyMuPDF: {e}")
    
    # Fallback to PyPDF2 (or use directly if PyMuPDF not available)
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text_content += page.extract_text() + "\n\n"
        return text_content.strip()
    except Exception as e2:
        print(f"Error extracting text from {file_path} with PyPDF2: {e2}")
        return ""

def create_text_chunks(text: str, chunk_size: int = 1200, overlap: int = 300) -> List[str]:
    """Split text into overlapping chunks for better retrieval with improved chunking."""
    if len(text) <= chunk_size:
        return [text]
    
    chunks = []
    start = 0
    
    while start < len(text):
        end = start + chunk_size
        
        # Try to break at a sentence, paragraph, or space boundary
        if end < len(text):
            # Look for the best break point within a reasonable range
            search_start = max(start + chunk_size // 2, start + 100)
            break_points = [
                text.rfind('.', search_start, end),
                text.rfind('!', search_start, end), 
                text.rfind('?', search_start, end),
                text.rfind('\n\n', search_start, end),
                text.rfind('\n', search_start, end),
                text.rfind('. ', search_start, end),
                text.rfind(' ', search_start, end)
            ]
            
            # Find the best break point
            best_break = max([bp for bp in break_points if bp > search_start], default=end)
            end = best_break + 1 if best_break != end else end
        
        chunk = text[start:end].strip()
        if chunk and len(chunk) > 50:  # Only include substantial chunks
            chunks.append(chunk)
        
        start = end - overlap if end < len(text) else end
    
    return chunks

'''
# Commented out - database/ directory removed
async def get_curated_qa_content() -> str:
    #Fetch curated Q&A content from database and format for knowledge base.
    try:
        if not AsyncSessionLocal:
            return ""
        async with AsyncSessionLocal() as db:
            from app.database.repositories import CuratedQaRepository
            curated_repo = CuratedQaRepository(db)
            content = await curated_repo.get_curated_qa_content_for_kb()
            return content
    except Exception as e:
        print(f"Error fetching curated Q&A content: {e}")
        return ""
'''

def is_mcl_document(file_path: Path) -> bool:
    """Determine if a document is MCL-related (not Spotplan)."""
    file_name_lower = file_path.name.lower()
    
    # Exclude Spotplan files explicitly
    spotplan_indicators = ['spotplan', 'spot_plan', 'spot-plan']
    if any(indicator in file_name_lower for indicator in spotplan_indicators):
        return False
    
    # Include MCL files
    mcl_indicators = [
        'mcl', 'checklist', 'creating', 'questions', 'knowledge', 'base',
        'rollenprofil', 'dashboard', 'tablet', 'phone', 'how-to-use',
        'aufgabe', 'mistakes', 'release', 'notes', 'quiz', 'dropbox',
        'business', 'case', 'tech_updates', 'vorgehen', 'app', 'tests'
    ]
    
    return any(indicator in file_name_lower for indicator in mcl_indicators)

def process_mcl_documents_with_enhanced_chunking() -> Tuple[List[str], List[Dict[str, Any]]]:
    """Process ONLY MCL documents with enhanced chunking and semantic search support."""
    global _mcl_document_chunks
    
    documents_path = Path("app/documents")
    file_ids = []
    _mcl_document_chunks = []
    chunk_id = 0
    
    print("\n" + "="*80)
    print("üöÄ PROCESSING MCL DOCUMENTS (ENHANCED VERSION)")
    print("="*80)
    
    # Initialize semantic search
    semantic_available = initialize_semantic_search()
    
    # Database content loading removed (database/ directory removed)
    # Previously loaded curated Q&A from database, now loads only from files
    
    # Process all file types in the documents directory
    if documents_path.exists():
        all_files = []
        supported_extensions = ['.pdf', '.docx', '.pptx', '.md']
        
        for ext in supported_extensions:
            all_files.extend(list(documents_path.glob(f"*{ext}")))
        
        # Filter to only MCL documents (exclude Spotplan)
        mcl_files = [f for f in all_files if is_mcl_document(f)]
        excluded_files = [f for f in all_files if not is_mcl_document(f)]
        
        print(f"\nüìÅ Found {len(all_files)} total files:")
        print(f"   ‚úÖ MCL files: {len(mcl_files)}")
        print(f"   ‚ùå Excluded files: {len(excluded_files)}")
        
        if excluded_files:
            print("   üö´ Excluded files (Spotplan/other):")
            for f in excluded_files:
                print(f"      - {f.name}")
        
        print(f"\nüìã Processing {len(mcl_files)} MCL files:")
        for mcl_file in mcl_files:
            print(f"   üìÑ {mcl_file.name}")
        
        for file_path in mcl_files:
            try:
                print(f"\nüîÑ Processing: {file_path.name}")
                
                # Extract text using appropriate method
                text_content = extract_text_from_file(file_path)
                
                if text_content.strip():
                    # Create chunks from the document
                    chunks = create_text_chunks(text_content)
                    
                    print(f"   üìä Created {len(chunks)} chunks")
                    
                    # Process each chunk
                    for i, chunk in enumerate(chunks):
                        # Create chunk metadata
                        chunk_metadata = {
                            "chunk_id": chunk_id,
                            "document_name": file_path.name,
                            "document_type": file_path.suffix.upper().replace('.', ''),
                            "chunk_index": i,
                            "total_chunks": len(chunks),
                            "content": chunk,
                            "content_hash": hashlib.md5(chunk.encode()).hexdigest()[:8],
                            "file_path": str(file_path)
                        }
                        _mcl_document_chunks.append(chunk_metadata)
                        
                        # Create a formatted chunk for OpenAI
                        formatted_chunk = f"""Document: {file_path.name}
Type: {chunk_metadata['document_type']}
Chunk {i+1}/{len(chunks)}

{chunk}

---
Source: {file_path.name} (Chunk {i+1})"""
                        
                        # Upload chunk to OpenAI
                        temp_file = BytesIO(formatted_chunk.encode('utf-8'))
                        created_file = client.files.create(
                            file=(f"{file_path.stem}_chunk_{i+1}.txt", temp_file),
                            purpose="assistants"
                        )
                        
                        file_ids.append(created_file.id)
                        chunk_id += 1
                    
                    print(f"   ‚úÖ Success: {len(chunks)} chunks created")
                else:
                    print(f"   ‚ö†Ô∏è No text content extracted from {file_path.name}")
                    
            except Exception as e:
                print(f"   ‚ùå Error processing {file_path.name}: {e}")
                traceback.print_exc()
    
    print(f"\nüìà SUMMARY:")
    print(f"   Total MCL chunks created: {len(_mcl_document_chunks)}")
    print(f"   OpenAI files uploaded: {len(file_ids)}")
    
    # Create embeddings if semantic search is available
    global _mcl_embeddings, _mcl_faiss_index
    if semantic_available and _mcl_document_chunks:
        print(f"   üîÑ Creating semantic embeddings...")
        _mcl_embeddings, _mcl_faiss_index = create_embeddings_for_chunks(_mcl_document_chunks)
        if _mcl_embeddings is not None:
            print(f"   ‚úÖ Semantic search ready with {len(_mcl_embeddings)} embeddings")
        else:
            print(f"   ‚ö†Ô∏è Semantic search setup failed")
    
    print("="*80 + "\n")
    return file_ids, _mcl_document_chunks

def expand_query_with_variants(query: str) -> List[str]:
    """Generate query variations to improve retrieval coverage."""
    variants = [query]
    
    # Add lowercase variant
    if query != query.lower():
        variants.append(query.lower())
    
    # Generate semantic variants using GPT
    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "system",
                    "content": "Generate 2-3 alternative phrasings of the user's question that capture the same intent but use different words. Return only the alternative questions, one per line, without numbering or explanation."
                },
                {"role": "user", "content": query}
            ],
            temperature=0.7,
            max_tokens=150
        )
        
        alternative_queries = response.choices[0].message.content.strip().split('\n')
        for alt_query in alternative_queries:
            clean_query = alt_query.strip().strip('123456789.-) ')
            if clean_query and len(clean_query) > 5:
                variants.append(clean_query)
        
        print(f"üìù Query expansion: {len(variants)} variants generated")
        for i, v in enumerate(variants[:4]):
            print(f"   Variant {i+1}: {v[:80]}")
        
    except Exception as e:
        print(f"‚ö†Ô∏è Query expansion failed: {e}")
    
    return variants[:4]  # Limit to 4 variants to avoid excessive searches

def find_relevant_chunks(query: str, max_chunks: int = 15) -> List[Dict[str, Any]]:
    """Find the most relevant document chunks using advanced RAG techniques.
    
    Implements:
    - Query expansion with variants
    - Hybrid search (semantic + keyword)
    - Translation for multilingual support
    - Increased context window (15 chunks vs 5)
    """
    print(f"\n[ADVANCED RAG] Starting search for query: '{query}'")
    
    # Detect language and translate if needed
    detected_lang = detect_language(query)
    search_query = query
    
    if detected_lang != 'en':
        print(f"[ADVANCED RAG] Detected non-English query ({detected_lang}), translating...")
        translated_query = translate_query_to_english(query, detected_lang)
        search_query = translated_query
        print(f"[ADVANCED RAG] Translated: '{search_query}'")
    
    # Generate query variants for better coverage
    query_variants = expand_query_with_variants(search_query)
    
    # Try semantic search first if available
    if HAS_SEMANTIC_SEARCH and _mcl_faiss_index is not None and _mcl_document_chunks:
        print(f"[ADVANCED RAG] Using semantic search with {len(query_variants)} query variants...")
        all_semantic_results = {}
        
        for variant in query_variants:
            semantic_results = semantic_search_chunks(
                variant, _mcl_document_chunks, _mcl_embeddings, _mcl_faiss_index, max_chunks * 2
            )
            # Aggregate results with scores
            for result in semantic_results:
                chunk_id = result['chunk_id']
                if chunk_id not in all_semantic_results:
                    all_semantic_results[chunk_id] = result
                else:
                    # Keep the highest score
                    if result.get('similarity_score', 0) > all_semantic_results[chunk_id].get('similarity_score', 0):
                        all_semantic_results[chunk_id] = result
        
        if all_semantic_results:
            # Sort by similarity score and return top results
            sorted_results = sorted(
                all_semantic_results.values(),
                key=lambda x: x.get('similarity_score', 0),
                reverse=True
            )[:max_chunks]
            
            print(f"[ADVANCED RAG] ‚úÖ Semantic search found {len(sorted_results)} unique high-quality matches")
            for i, r in enumerate(sorted_results[:5]):
                print(f"   #{i+1}: {r['document_name']} (Score: {r.get('similarity_score', 0):.3f})")
            
            return sorted_results
        else:
            print(f"[ADVANCED RAG] ‚ö†Ô∏è No semantic matches, falling back to enhanced keyword search")
    
    # Fallback to enhanced keyword search with query expansion
    print(f"[ADVANCED RAG] Using enhanced keyword search...")
    query_lower = search_query.lower()
    print(f"[ADVANCED RAG] Primary search query: '{query_lower}'")
    
    # Enhanced keyword matching with query expansion
    chunk_scores = {}  # Use dict to aggregate scores across variants
    
    for variant_query in query_variants:
        variant_lower = variant_query.lower()
        query_words = variant_lower.split()
        
        # Generate 2-word and 3-word phrases
        phrases = [variant_lower]
        for i in range(len(query_words) - 1):
            phrases.append(" ".join(query_words[i:i+2]))
            if i < len(query_words) - 2:
                phrases.append(" ".join(query_words[i:i+3]))
        
        for chunk in _mcl_document_chunks:
            chunk_id = chunk['chunk_id']
            content_lower = chunk["content"].lower()
            doc_name_lower = chunk["document_name"].lower()
            
            # Initialize score for this chunk if not exists
            if chunk_id not in chunk_scores:
                chunk_scores[chunk_id] = {'chunk': chunk, 'score': 0, 'matches': set()}
            
            # Score individual words (weight: 1)
            for word in query_words:
                if len(word) > 2:
                    word_count = content_lower.count(word)
                    if word_count > 0:
                        chunk_scores[chunk_id]['score'] += word_count
                        chunk_scores[chunk_id]['matches'].add(word)
            
            # Score phrases (weight: 5)
            for phrase in phrases:
                if len(phrase) > 4:
                    phrase_count = content_lower.count(phrase)
                    if phrase_count > 0:
                        chunk_scores[chunk_id]['score'] += phrase_count * 5
                        chunk_scores[chunk_id]['matches'].add(f'[{phrase}]')
            
            # Document type bonuses
            if any(term in doc_name_lower for term in ["how-to", "guide", "creating", "tutorial"]):
                chunk_scores[chunk_id]['score'] += 3
            
            # Specific content bonuses
            if "task" in variant_lower and "task" in doc_name_lower:
                chunk_scores[chunk_id]['score'] += 5
            
            if "question" in variant_lower and "question" in doc_name_lower:
                chunk_scores[chunk_id]['score'] += 5
            
            if any(term in variant_lower for term in ["create", "creat", "new"]) and \
               "creating" in doc_name_lower:
                chunk_scores[chunk_id]['score'] += 8
            
            if "login" in variant_lower and any(term in content_lower for term in ["login", "benutzername", "password"]):
                chunk_scores[chunk_id]['score'] += 10
    
    # Sort by score and return top chunks
    sorted_chunks = sorted(
        chunk_scores.values(),
        key=lambda x: x['score'],
        reverse=True
    )
    
    top_chunks = [item['chunk'] for item in sorted_chunks[:max_chunks] if item['score'] > 0]
    
    print(f"[ADVANCED RAG] Found {len([s for s in chunk_scores.values() if s['score'] > 0])} relevant chunks")
    print(f"[ADVANCED RAG] Returning top {len(top_chunks)} chunks:")
    for i, item in enumerate(sorted_chunks[:5]):
        chunk = item['chunk']
        print(f"   #{i+1}: {chunk['document_name']} (Score: {item['score']})")
    
    if len(top_chunks) == 0:
        print(f"[ADVANCED RAG] ‚ö†Ô∏è WARNING: No relevant chunks found!")
    
    # Hybrid search: Combine semantic and keyword results with re-ranking
    semantic_chunks = []
    if len(top_chunks) < max_chunks:
        print(f"[ADVANCED RAG] Enhancing with semantic search...")
        # Use semantic search to find additional relevant chunks
        semantic_chunks = semantic_search_chunks(query, max_results=max_chunks)
        print(f"[ADVANCED RAG] Semantic search found {len(semantic_chunks)} chunks")
    
    # Merge and re-rank results
    final_chunks = []
    seen_chunk_ids = set()
    
    # First, add keyword-scored chunks (already sorted by score)
    for chunk in top_chunks:
        if chunk['chunk_id'] not in seen_chunk_ids:
            final_chunks.append({
                'chunk': chunk,
                'keyword_score': chunk_scores[chunk['chunk_id']]['score'],
                'semantic_score': 0,
                'source': 'keyword'
            })
            seen_chunk_ids.add(chunk['chunk_id'])
    
    # Then add semantic chunks with their scores
    for semantic_chunk in semantic_chunks:
        chunk_id = semantic_chunk['chunk']['chunk_id']
        if chunk_id not in seen_chunk_ids:
            final_chunks.append({
                'chunk': semantic_chunk['chunk'],
                'keyword_score': chunk_scores.get(chunk_id, {}).get('score', 0),
                'semantic_score': semantic_chunk['score'],
                'source': 'semantic'
            })
            seen_chunk_ids.add(chunk_id)
        else:
            # Update existing entry with semantic score
            for item in final_chunks:
                if item['chunk']['chunk_id'] == chunk_id:
                    item['semantic_score'] = semantic_chunk['score']
                    item['source'] = 'hybrid'
                    break
    
    # Re-rank using combined score (60% keyword, 40% semantic)
    for item in final_chunks:
        item['combined_score'] = (0.6 * item['keyword_score']) + (0.4 * item['semantic_score'] * 100)
    
    final_chunks.sort(key=lambda x: x['combined_score'], reverse=True)
    
    # Return top max_chunks
    result = [item['chunk'] for item in final_chunks[:max_chunks]]
    
    print(f"[ADVANCED RAG] Final re-ranked results ({len(result)} chunks):")
    for i, item in enumerate(final_chunks[:5]):
        chunk = item['chunk']
        print(f"   #{i+1}: {chunk['document_name']} [Source: {item['source']}, Combined: {item['combined_score']:.1f}]")
    
    return result

def start_mcl_knowledge_base() -> str:
    """Initialize the MCL knowledge base with all documents."""
    global _mcl_vector_store_id
    
    print("Starting MCL knowledge base initialization...")
    try:
        # Process all MCL documents with chunk tracking
        file_ids, chunks = process_mcl_documents_with_enhanced_chunking()
        
        if not file_ids:
            print("WARNING: No files were processed for the MCL knowledge base")
            return None
        
        print(f"Creating vector store with {len(file_ids)} file chunks...")
        vector_store = client.vector_stores.create(
            name="mcl_knowledge_base_chunked",
            file_ids=file_ids
        )
        
        if not vector_store.id:
            raise ValueError("Failed to create a vector store with a valid ID.")

        _mcl_vector_store_id = vector_store.id
        print(f"MCL knowledge base setup complete. Vector Store ID: {vector_store.id}")
        print(f"Total document chunks indexed: {len(chunks)}")
        
        return vector_store.id

    except Exception as e:
        print(f"FATAL: An error occurred during MCL knowledge base setup: {e}")
        return None

def detect_language(text: str) -> str:
    """Detect the language of the user's query using GPT-4.
    Returns language code: 'de' for German, 'en' for English, etc.
    """
    if not text or len(text.strip()) < 2:
        return 'en'  # Default to English
    
    try:
        # Quick heuristic check first (fast path)
        text_lower = text.lower()
        
        # Check for German-specific characters
        if any(char in text_lower for char in ['√§', '√∂', '√º', '√ü']):
            return 'de'
        
        # Check for obvious German words
        strong_german_words = ['ich', 'kannst', 'mir', 'checkliste', 'erstellen', 'anlegen', 
                               'wie', 'was', 'warum', 'erkl√§ren', 'hilfe', 'bitte']
        if any(word in text_lower for word in strong_german_words):
            return 'de'
        
        # If still uncertain and text is long enough, use GPT for detection
        if len(text.split()) > 3:
            response = client.chat.completions.create(
                model="gpt-4o-mini",  # Faster and cheaper for language detection
                messages=[
                    {"role": "system", "content": "Detect the language of the following text. Respond with ONLY the ISO 639-1 language code (e.g., 'en' for English, 'de' for German, 'es' for Spanish, etc.). No explanation."},
                    {"role": "user", "content": text}
                ],
                temperature=0,
                max_tokens=10
            )
            detected_lang = response.choices[0].message.content.strip().lower()
            return detected_lang if detected_lang in ['de', 'en', 'es', 'fr', 'it'] else 'en'
        
        return 'en'  # Default to English
        
    except Exception as e:
        print(f"Error detecting language: {e}")
        return 'en'  # Default to English on error

def translate_query_to_english(query: str, source_language: str) -> str:
    """Translate a user query to English for document search.
    
    Args:
        query: The user's query in any language
        source_language: The detected source language code
    
    Returns:
        English translation of the query
    """
    if source_language == 'en':
        return query  # Already in English
    
    try:
        print(f"üîÑ Translating query from {source_language} to English...")
        
        response = client.chat.completions.create(
            model="gpt-4o-mini",  # Fast and accurate for translation
            messages=[
                {
                    "role": "system", 
                    "content": "You are a professional translator. Translate the following text to English. Preserve the meaning and intent. Respond with ONLY the English translation, no explanations or additional text."
                },
                {
                    "role": "user", 
                    "content": f"Translate to English: {query}"
                }
            ],
            temperature=0.3,
            max_tokens=200
        )
        
        translated = response.choices[0].message.content.strip()
        print(f"‚úÖ Translation: '{query}' ‚Üí '{translated}'")
        return translated
        
    except Exception as e:
        print(f"‚ùå Error translating query: {e}")
        return query  # Return original on error

def get_mcl_ai_response(messages_input: List[Dict[str, Any]]) -> Any:
    """Get AI response for MCL-related queries with source attribution."""
    
    print("\n" + "="*80)
    print("[MCL AI] Starting MCL AI Response Generation")
    print("="*80)
    
    # Extract the latest user message for relevance search
    latest_user_message = ""
    for msg in reversed(messages_input):
        if msg.get("role") == "user":
            latest_user_message = msg.get("content", "")
            break
    
    print(f"[MCL AI] User message: '{latest_user_message}'")
    
    # Find relevant document chunks using enhanced search
    relevant_chunks = find_relevant_chunks(latest_user_message, max_chunks=5)
    
    # Create context from relevant chunks
    context_parts = []
    sources = []
    
    print(f"\n[MCL AI] Building context from {len(relevant_chunks)} relevant chunks:")
    for i, chunk in enumerate(relevant_chunks):
        print(f"[MCL AI] Context Chunk {i+1}: {chunk['document_name']} (Chunk {chunk['chunk_index']+1}/{chunk['total_chunks']})")
        print(f"[MCL AI] Preview: {chunk['content'][:150]}...")
        
        context_parts.append(f"[From {chunk['document_name']}, Chunk {chunk['chunk_index']+1}]:\n{chunk['content']}")
        source_info = f"{chunk['document_name']} (Chunk {chunk['chunk_index']+1}/{chunk['total_chunks']})"
        if source_info not in sources:
            sources.append(source_info)
    
    context = "\n\n" + "\n\n---\n\n".join(context_parts) if context_parts else ""
    
    if not context:
        print("[MCL AI WARNING] No context available! AI will have no document excerpts to work with.")
    else:
        print(f"[MCL AI] Total context length: {len(context)} characters")
    
    # Detect user's language
    detected_lang = detect_language(latest_user_message)
    user_language_name = {
        'de': 'German',
        'en': 'English',
        'es': 'Spanish',
        'fr': 'French',
        'it': 'Italian'
    }.get(detected_lang, 'English')
    
    print(f"\n[MCL AI] Detected language: {user_language_name} ({detected_lang})")
    
    # Create language-specific system prompt with enhanced reasoning capabilities
    if detected_lang == 'de':
        print("[MCL AI] Creating German-language system prompt")
        system_prompt = f"""üá©üá™ WICHTIG: Du musst auf DEUTSCH antworten! Der Benutzer stellt eine Frage auf Deutsch.

You are "MCL Assistant," an expert AI assistant for the MCL (Mobile Checklist) application.

‚ö†Ô∏è CRITICAL LANGUAGE RULE: The user is writing in GERMAN. You MUST respond in GERMAN (Deutsch).
- Translate ALL information from the English documents into clear, natural German
- Use proper German terminology:
  * "Checklist" ‚Üí "Checkliste" or "Pr√ºfliste"
  * "Task" ‚Üí "Aufgabe"
  * "Question" ‚Üí "Frage"
  * "Dashboard" ‚Üí "Dashboard" (same)
  * "Create" ‚Üí "erstellen" or "anlegen"
  * "User" ‚Üí "Benutzer"
  * "Wizard" ‚Üí "Assistent"
  * "Click" ‚Üí "Klicken"
  * "Button" ‚Üí "Schaltfl√§che" or "Button"
- Write complete sentences in German with proper grammar
- Provide step-by-step instructions in German

Available Document Excerpts (in English - you must translate):
{context}

Guidelines:
- Use the provided document excerpts as your PRIMARY source of information
- You can synthesize information across multiple excerpts to provide comprehensive answers
- If information spans multiple documents, combine them into a coherent response
- Always cite which document(s) you're referencing (in German: "Quelle:" or "Aus:")
- If the user's question requires information not explicitly stated but logically implied, you may provide reasonable inferences while noting they are inferences
- If critical information is missing, say in German: "Diese spezifische Information finde ich nicht in den verf√ºgbaren Dokumenten."
- Provide step-by-step instructions in German when available
- Be specific and detailed based on the actual documentation
- At the end of your response, list the sources in German: "üìö Quellen:"

REMEMBER: Your ENTIRE response must be in GERMAN (Deutsch)!"""
    else:
        print("[MCL AI] Creating English-language system prompt")
        system_prompt = f"""You are "MCL Assistant," an expert AI assistant for the MCL (Mobile Checklist) application. 

Your goal is to provide helpful, accurate answers based on the MCL documentation provided below.

Available Document Excerpts:
{context}

Guidelines:
- Use these document excerpts as your PRIMARY source of information
- You can synthesize information across multiple excerpts to provide comprehensive answers
- If information spans multiple documents, combine them into a coherent response
- Always cite which document(s) you're referencing (e.g., "According to the Creating Checklists guide...")
- If the user's question requires information not explicitly stated but logically implied, you may provide reasonable inferences while clearly noting they are inferences
- If critical information is clearly missing, say: "I don't have specific information about [X] in the available documents."
- Provide clear, step-by-step instructions when available in the documents
- Be specific and detailed based on the actual documentation
- At the end of your response, list the sources you used: "üìö Sources:"

Special topics to be aware of:
- For questions about creating tasks, checklists, or questions: refer to the "Creating" guides
- For questions about using the app: refer to "How-to-use" guides for Phone, Tablet, or Dashboard
- For login and access: refer to the How-to-use guides (login is usually the first step)
- For question types: refer to "Creating Questions" document

Remember: Base your answers primarily on the provided excerpts, but you can make logical connections between information from different sections."""
    
    final_messages = [{"role": "system", "content": system_prompt}] + messages_input
    
    print(f"\n[MCL AI] System prompt length: {len(system_prompt)} characters")
    print(f"[MCL AI] Total messages to send: {len(final_messages)}")
    print(f"[MCL AI] Sending request to OpenAI (Language: {user_language_name})...")

    try:
        response = client.chat.completions.create( 
            model="gpt-4o",
            messages=final_messages,
            temperature=0.2,  # Slightly higher for more natural translations
            max_tokens=2000
        )
        
        print(f"[MCL AI] Received response from OpenAI")
        
        # Enhance the response with source information
        original_content = response.choices[0].message.content
        
        print(f"[MCL AI] Response content length: {len(original_content)} characters")
        print(f"[MCL AI] Response preview: {original_content[:200]}...")
        
        if sources:
            # Use language-appropriate source header
            source_headers = {
                'de': "\n\nüìö **Quellen:**\n",
                'en': "\n\nüìö **Sources:**\n",
                'es': "\n\nüìö **Fuentes:**\n",
                'fr': "\n\nüìö **Sources:**\n"
            }
            sources_header = source_headers.get(detected_lang, source_headers['en'])
            sources_text = sources_header + "\n".join([f"‚Ä¢ {source}" for source in sources])
            
            enhanced_content = original_content + sources_text
            
            # Create enhanced response
            response.choices[0].message.content = enhanced_content
        
        print(f"[MCL AI] ‚úì Response generation completed successfully")
        print(f"[MCL AI] Sources attached: {len(sources)}")
        print("="*80 + "\n")
        
        return response
        
    except Exception as e:
        print(f"[MCL AI ERROR] Exception occurred: {e}")
        import traceback
        traceback.print_exc()
        print("="*80 + "\n")
        
        # Create a fallback response
        class MockResponse:
            def __init__(self, content):
                self.choices = [MockChoice(content)]
        
        class MockChoice:
            def __init__(self, content):
                self.message = MockMessage(content)
        
        class MockMessage:
            def __init__(self, content):
                self.role = "assistant"
                self.content = content
                self.tool_calls = None
        
        # Create language-appropriate fallback message
        fallback_messages = {
            'de': """Entschuldigung, aber ich habe derzeit technische Schwierigkeiten beim Zugriff auf die MCL-Wissensdatenbank.
            
            Bitte stellen Sie Ihre Frage erneut, oder wenden Sie sich f√ºr technischen Support an Ihren Systemadministrator.""",
            'en': """I apologize, but I'm currently experiencing technical difficulties accessing the MCL knowledge base.
            
            Please try your question again, or contact your system administrator for technical support.""",
            'es': """Lo siento, pero actualmente tengo dificultades t√©cnicas para acceder a la base de conocimientos de MCL.
            
            Por favor, intente su pregunta nuevamente o contacte a su administrador del sistema para soporte t√©cnico."""
        }
        
        fallback_content = fallback_messages.get(detected_lang, fallback_messages['en'])
        return MockResponse(fallback_content)

def debug_mcl_knowledge_base():
    """Debug function to show what's in the MCL knowledge base."""
    print("\n" + "="*80)
    print("üîç MCL KNOWLEDGE BASE DEBUG REPORT")
    print("="*80)
    
    print(f"üìä Total MCL chunks: {len(_mcl_document_chunks)}")
    print(f"üß† Semantic search available: {HAS_SEMANTIC_SEARCH}")
    print(f"üìö Embeddings created: {_mcl_embeddings is not None}")
    print(f"üîé FAISS index ready: {_mcl_faiss_index is not None}")
    
    if _mcl_document_chunks:
        # Group by document
        docs = {}
        for chunk in _mcl_document_chunks:
            doc_name = chunk['document_name']
            if doc_name not in docs:
                docs[doc_name] = []
            docs[doc_name].append(chunk)
        
        print(f"\nüìÅ Documents in MCL knowledge base:")
        for doc_name, chunks in docs.items():
            print(f"   üìÑ {doc_name}: {len(chunks)} chunks ({chunks[0]['document_type']})")
        
        print(f"\nüîç Sample chunks (first 3):")
        for i, chunk in enumerate(_mcl_document_chunks[:3]):
            print(f"   Chunk {i+1}: {chunk['document_name']} - {chunk['content'][:100]}...")
    else:
        print("‚ö†Ô∏è No chunks found in MCL knowledge base!")
    
    print("="*80 + "\n")

def get_ai_format(messages_input, request: BaseModel):
    final_messages = [messages_input]

    print(f"Sending messages to AI for function-calling: {final_messages}")

    response = client.beta.chat.completions.parse(
        model="gpt-4o",
        messages=final_messages,
        response_format=request.model_dump()
    )
    print(f"Received AI response object: {type(response)}")
    return response


# ========================================================================
# Vision-Enabled Chat Support
# ========================================================================

def get_vision_enabled_response(messages: list, vector_store_id: str) -> dict:
    """
    Handle chat messages with optional image attachments using GPT-4o vision.
    
    This function processes messages that may contain both text and images,
    similar to ChatGPT/Gemini interfaces. It:
    1. Detects if any message contains images
    2. If images present, uses vision-enabled response
    3. If no images, falls back to regular RAG response
    4. Supports multimodal conversations
    
    Args:
        messages: List of message dicts with role and content (may include images)
        vector_store_id: MCL knowledge base vector store ID
    
    Returns:
        dict with 'response', 'success', 'has_vision', and optional 'metadata'
    """
    print("\n" + "="*80)
    print("üîç VISION-ENABLED CHAT ANALYSIS")
    print("="*80)
    
    # Check if any message contains images
    has_images = False
    image_messages = []
    
    for msg in messages:
        if isinstance(msg.get('content'), list):
            for item in msg['content']:
                if isinstance(item, dict) and item.get('type') == 'image_url':
                    has_images = True
                    image_messages.append(msg)
                    print(f"üì∑ Image detected in {msg.get('role', 'unknown')} message")
    
    if not has_images:
        print("‚ÑπÔ∏è No images detected - using standard RAG response")
        print("="*80 + "\n")
        # Fall back to regular MCL AI response
        response_obj = get_mcl_ai_response(messages)
        return {
            "response": response_obj.choices[0].message.content,
            "success": True,
            "has_vision": False
        }
    
    print(f"üñºÔ∏è Vision mode activated - {len(image_messages)} message(s) with images")
    
    try:
        # Extract the latest user query and images
        latest_user_message = None
        for msg in reversed(messages):
            if msg.get('role') == 'user':
                latest_user_message = msg
                break
        
        if not latest_user_message:
            return {
                "response": "No user message found",
                "success": False,
                "has_vision": True,
                "error": "No user message in conversation"
            }
        
        # Extract text query from multimodal content
        user_query = ""
        image_data_list = []
        
        if isinstance(latest_user_message.get('content'), list):
            for item in latest_user_message['content']:
                if isinstance(item, dict):
                    if item.get('type') == 'text':
                        user_query = item.get('text', '')
                    elif item.get('type') == 'image_url':
                        img_url = item.get('image_url', {}).get('url', '')
                        if img_url:
                            image_data_list.append(img_url)
        elif isinstance(latest_user_message.get('content'), str):
            user_query = latest_user_message['content']
        
        print(f"üí¨ User query: {user_query[:100]}..." if len(user_query) > 100 else f"üí¨ User query: {user_query}")
        print(f"üñºÔ∏è Images to analyze: {len(image_data_list)}")
        
        # Build vision-enabled messages for OpenAI
        vision_messages = []
        
        # Add system context from MCL knowledge base
        system_prompt = """You are the MCL App Assistant with vision capabilities. When analyzing screenshots:

1. Identify if the image is from the MCL (Mobile Checklist) App
2. If yes, describe what screen/feature is shown
3. Answer the user's question based on both the visual context and your MCL knowledge
4. If no, politely inform the user you can only help with MCL App screenshots
5. Provide clear, step-by-step guidance when needed

Be helpful, concise, and accurate."""
        
        vision_messages.append({
            "role": "system",
            "content": system_prompt
        })
        
        # Add conversation history (excluding images from history to save tokens)
        for msg in messages[:-1]:  # All except the latest
            if msg.get('role') in ['user', 'assistant', 'system']:
                content = msg.get('content', '')
                # Extract text only from multimodal messages
                if isinstance(content, list):
                    text_parts = [item.get('text', '') for item in content if item.get('type') == 'text']
                    content = ' '.join(text_parts)
                
                vision_messages.append({
                    "role": msg['role'],
                    "content": content
                })
        
        # Add the latest user message with images
        multimodal_content = []
        
        if user_query:
            multimodal_content.append({
                "type": "text",
                "text": user_query
            })
        
        for img_data in image_data_list:
            multimodal_content.append({
                "type": "image_url",
                "image_url": {
                    "url": img_data,
                    "detail": "high"  # Use high detail for better analysis
                }
            })
        
        vision_messages.append({
            "role": "user",
            "content": multimodal_content
        })
        
        print(f"üì® Sending {len(vision_messages)} messages to GPT-4o vision...")
        
        # Call OpenAI with vision support
        response = client.chat.completions.create(
            model="gpt-4o",  # GPT-4o supports vision
            messages=vision_messages,
            max_tokens=1500,
            temperature=0.7
        )
        
        ai_response = response.choices[0].message.content
        
        print(f"‚úÖ Vision response received ({len(ai_response)} characters)")
        print(f"üìù Response preview: {ai_response[:150]}..." if len(ai_response) > 150 else f"üìù Response: {ai_response}")
        print("="*80 + "\n")
        
        return {
            "response": ai_response,
            "success": True,
            "has_vision": True,
            "metadata": {
                "model": "gpt-4o",
                "images_analyzed": len(image_data_list),
                "query": user_query
            }
        }
        
    except Exception as e:
        print(f"‚ùå ERROR in vision processing: {e}")
        import traceback
        traceback.print_exc()
        print("="*80 + "\n")
        
        return {
            "response": f"I encountered an error processing the image: {str(e)}",
            "success": False,
            "has_vision": True,
            "error": str(e)
        }

